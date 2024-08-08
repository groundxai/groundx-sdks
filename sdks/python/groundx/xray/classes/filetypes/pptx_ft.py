# coding: utf-8


from typing import Tuple

import pptx
from pptx.shapes.autoshape import Shape
from pptx.shapes.placeholder import SlidePlaceholder

from groundx.xray.classes.document import XRayDocument
from groundx.xray.classes.filetypes.constants import TXT_MODIFIER


class XRayPPTX(XRayDocument):
    def __init__(self, file_path: str) -> None:
        super().__init__(file_path)

        self.pages, self.tokens = self.estimate_size()

    def _estimate_size(self) -> Tuple[int, int]:
        prs = pptx.Presentation(self.file_path)

        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if isinstance(shape, (Shape, SlidePlaceholder)):
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            full_text.append(paragraph.text)

        return len(prs.slides), self._token_count('\n'.join(full_text)) * TXT_MODIFIER